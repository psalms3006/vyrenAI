"""Universal URL Understanding subsystem for VYREN.

This module detects, classifies, resolves, and extracts content from URLs,
then normalizes the result into a common representation that downstream
VYREN systems can use without caring whether the source was Instagram,
YouTube, a PDF, GitHub, or a plain HTML page.

Design goals:
  - reuse existing VYREN subsystems (browser, vision, OCR, document tools)
  - avoid hardcoded per-site scraping hacks
  - graceful failure with explicit status
  - session-scoped conversational context by default
  - optional memory/knowledge-graph integration
"""
from __future__ import annotations

import json
import os
import re
import tempfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional
from urllib.parse import urlparse

# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------

class ExtractionQuality(str):
    FULL = "full"
    PARTIAL = "partial"
    METADATA_ONLY = "metadata_only"
    FAILED = "failed"

class ResourceStatus(str):
    SUCCESS = "success"
    PARTIAL = "partial"
    INACCESSIBLE = "inaccessible"
    AUTH_REQUIRED = "auth_required"
    BLOCKED = "blocked"
    UNSUPPORTED = "unsupported"
    RATE_LIMITED = "rate_limited"
    MALFORMED_URL = "malformed_url"
    NETWORK_FAILURE = "network_failure"
    UNKNOWN = "unknown"


@dataclass
class UrlResource:
    url: str
    canonical_url: str = ""
    source_type: str = "unknown"
    platform: str = ""
    title: str = ""
    author: str = ""
    published_at: str = ""
    description: str = ""
    text: str = ""
    markdown: str = ""
    media: list[str] = field(default_factory=list)
    images: list[str] = field(default_factory=list)
    videos: list[str] = field(default_factory=list)
    documents: list[str] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    extraction_method: str = ""
    extraction_quality: str = ExtractionQuality.FAILED
    retrieved_at: str = field(default_factory=lambda: datetime.now(timezone.utc).isoformat())
    status: str = ResourceStatus.UNKNOWN
    errors: list[str] = field(default_factory=list)
    access_state: str = ""

    def to_dict(self) -> dict:
        return {
            "url": self.url,
            "canonical_url": self.canonical_url,
            "source_type": self.source_type,
            "platform": self.platform,
            "title": self.title,
            "author": self.author,
            "published_at": self.published_at,
            "description": self.description,
            "text": self.text[:4000],
            "markdown": self.markdown[:4000],
            "media": self.media,
            "images": self.images,
            "videos": self.videos,
            "documents": self.documents,
            "metadata": self.metadata,
            "extraction_method": self.extraction_method,
            "extraction_quality": self.extraction_quality,
            "retrieved_at": self.retrieved_at,
            "status": self.status,
            "errors": self.errors,
            "access_state": self.access_state,
        }

    def to_conversation_context(self, max_chars: int = 3000) -> str:
        parts = []
        if self.title:
            parts.append(f"Title: {self.title}")
        if self.author:
            parts.append(f"Author: {self.author}")
        if self.published_at:
            parts.append(f"Published: {self.published_at}")
        if self.description:
            parts.append(f"Description: {self.description}")
        if self.text:
            trimmed = self.text[:max_chars]
            parts.append(f"Content:\n{trimmed}")
        if self.markdown:
            trimmed = self.markdown[:max_chars]
            parts.append(f"Markdown:\n{trimmed}")
        if self.status not in (ResourceStatus.SUCCESS, ResourceStatus.PARTIAL):
            parts.append(f"Note: extraction status={self.status}; quality={self.extraction_quality}")
        if self.errors:
            parts.append("Errors:\n- " + "\n- ".join(self.errors[:5]))
        return "\n".join(parts) if parts else ""

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_URL_RE = re.compile(
    r"https?://(?:[A-Za-z0-9\-._~:/?#\[\]@!$&'()*+,;=%]+?)(?:[^A-Za-z0-9\-._~:/?#\[\]@!$&'()*+,;=%]|$)",
    re.IGNORECASE,
)

_PLATFORM_PATTERNS: dict[str, list[str]] = {
    "youtube": ["youtube.com", "youtu.be"],
    "twitter": ["x.com", "twitter.com"],
    "reddit": ["reddit.com"],
    "github": ["github.com"],
    "instagram": ["instagram.com"],
    "linkedin": ["linkedin.com"],
    "tiktok": ["tiktok.com"],
    "facebook": ["facebook.com"],
    "docs": ["docs.google.com", "notion.so"],
    "pdf": [".pdf"],
    "docx": [".docx"],
    "xlsx": [".xlsx", ".xls"],
    "pptx": [".pptx"],
}

_EXTRA_TYPE_PATTERNS: dict[str, list[str]] = {
    "video": [".mp4", ".mov", ".m3u8", ".webm"],
    "audio": [".mp3", ".wav", ".ogg", ".m4a"],
    "csv": [".csv"],
    "markdown": [".md", ".markdown"],
    "json": [".json"],
}


def detect_urls(text: str) -> list[str]:
    urls = []
    for match in _URL_RE.finditer(text):
        url = match.group(0)
        url = url.rstrip(".,;:!?)]>'\"")
        if url:
            urls.append(url)
    return list(dict.fromkeys(urls))


def classify_url(url: str) -> dict[str, str]:
    parsed = urlparse(url)
    host = parsed.netloc.lower()
    path = parsed.path.lower()
    query = (parsed.query or "").lower()
    platform = ""
    source_type = "web"

    for plat, patterns in _PLATFORM_PATTERNS.items():
        for pat in patterns:
            if pat.startswith("."):
                if path.endswith(pat):
                    platform = plat if not platform else platform
                    source_type = "document" if plat in {"pdf", "docx", "xlsx", "pptx"} else source_type
            elif pat in host or pat in path or pat in query:
                platform = plat
                break

    if not platform:
        if any(path.endswith(ext) for ext in [".pdf", ".docx", ".xlsx", ".pptx", ".md", ".csv", ".json"]):
            source_type = "document"
    if platform in {"pdf", "docx", "xlsx", "pptx", "csv", "markdown", "json"}:
        source_type = "document"

    if platform in {"video", "audio"}:
        source_type = "media"

    return {"url": url, "source_type": source_type, "platform": platform or "web"}

# ---------------------------------------------------------------------------
# Extraction subsystem
# ---------------------------------------------------------------------------

class UrlExtractor:
    """Reuses existing VYREN subsystems instead of introducing new deps."""

    def __init__(self, registry: Any | None = None, tmp_root: Path | None = None) -> None:
        self._registry = registry
        self._tmp_root = Path(tmp_root or Path(tempfile.gettempdir()) / "vyren_urls")

    def extract(self, url: str) -> UrlResource:
        resource = UrlResource(url=url)
        classification = classify_url(url)
        resource.source_type = classification["source_type"]
        resource.platform = classification["platform"]

        if not _URL_RE.match(url):
            resource.status = ResourceStatus.MALFORMED_URL
            resource.errors.append("URL does not match http/https pattern")
            return resource

        # Direct HTTP extraction first for lightweight pages/docs
        try:
            http_result = _try_http_extraction(url)
        except Exception as exc:
            http_result = None
            resource.errors.append(f"http_extraction_error={exc}")

        if http_result and http_result.get("text") or http_result and http_result.get("markdown"):
            resource.title = http_result.get("title", "")
            resource.description = http_result.get("description", "")
            resource.text = http_result.get("text", "")
            resource.markdown = http_result.get("markdown", "")
            resource.metadata.update(http_result.get("metadata", {}))
            resource.extraction_method = "http"
            resource.status = ResourceStatus.SUCCESS
            resource.extraction_quality = ExtractionQuality.FULL
            resource.canonical_url = http_result.get("canonical_url", url)
            return resource

        # Browser fallback for JS pages and weak HTTP extraction
        browser_result = _try_browser_extraction(url)
        if browser_result:
            resource.title = browser_result.get("title", "")
            resource.description = browser_result.get("description", "")
            resource.text = browser_result.get("text", "")
            resource.markdown = browser_result.get("markdown", "")
            resource.images = browser_result.get("images", [])
            resource.videos = browser_result.get("videos", [])
            resource.metadata.update(browser_result.get("metadata", {}))
            resource.extraction_method = "browser"
            resource.status = ResourceStatus.SUCCESS if resource.text or resource.markdown else ResourceStatus.PARTIAL
            resource.extraction_quality = ExtractionQuality.FULL if resource.text else ExtractionQuality.METADATA_ONLY
            resource.canonical_url = browser_result.get("canonical_url", url)
            return resource

        # Document fallback for direct document URLs
        if resource.source_type == "document":
            doc_result = _try_document_extraction(url)
            if doc_result:
                resource.text = doc_result.get("text", "")
                resource.markdown = doc_result.get("markdown", "")
                resource.metadata.update(doc_result.get("metadata", {}))
                resource.extraction_method = "document"
                resource.status = ResourceStatus.SUCCESS if resource.text else ResourceStatus.PARTIAL
                resource.extraction_quality = ExtractionQuality.FULL if resource.text else ExtractionQuality.METADATA_ONLY
                return resource

        # Last-resort metadata probe
        meta = _try_metadata_only(url)
        if meta:
            resource.title = meta.get("title", "")
            resource.description = meta.get("description", "")
            resource.metadata.update(meta.get("metadata", {}))
            resource.extraction_method = "metadata"
            resource.status = ResourceStatus.PARTIAL
            resource.extraction_quality = ExtractionQuality.METADATA_ONLY
            return resource

        resource.status = ResourceStatus.INACCESSIBLE
        resource.errors.append("No extraction path produced readable content")
        return resource

# ---------------------------------------------------------------------------
# Extraction implementations
# ---------------------------------------------------------------------------

def _try_http_extraction(url: str) -> dict | None:
    parsed = urlparse(url)
    path = parsed.path.lower()
    query = (parsed.query or "").lower()

    if path.endswith(".pdf") or "pdf" in query:
        return None
    if path.endswith((".docx", ".doc")) or "docx" in query or "document" in query:
        return None
    if path.endswith((".xlsx", ".xls")) or "spreadsheets" in query:
        return None
    if path.endswith((".pptx",)) or "presentation" in query:
        return None

    try:
        import httpx
    except Exception:
        return None

    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
    }
    try:
        resp = httpx.get(url, headers=headers, follow_redirects=True, timeout=12)
    except Exception:
        return None

    if resp.status_code >= 400:
        return None
    if resp.status_code == 403:
        return None

    content_type = resp.headers.get("content-type", "")
    text = resp.text or ""
    if not text.strip():
        return None

    title = _extract_title(text)
    description = _extract_meta(text, "description")
    canonical = _extract_canonical(text, url)
    body = _extract_readable_text(text)
    markdown = _html_to_markdown(text)
    metadata = {
        "content_type": content_type,
        "final_url": str(resp.url),
        "status_code": resp.status_code,
    }

    if not body and not markdown:
        return None

    return {
        "title": title,
        "description": description,
        "text": body or "",
        "markdown": markdown or "",
        "canonical_url": canonical or str(resp.url),
        "metadata": metadata,
    }


def _try_browser_extraction(url: str) -> dict | None:
    try:
        from browser import go_to, get_text, screenshot
    except Exception:
        return None

    try:
        go_to(url=url, timeout=20)
    except Exception:
        return None

    page_text = ""
    try:
        page_text = get_text(timeout=20) or ""
    except Exception:
        page_text = ""

    title = ""
    try:
        from browser import _get_title
        title = _get_title() or ""
    except Exception:
        title = ""

    md = _text_to_markdown(page_text)
    result = {
        "title": title,
        "text": page_text,
        "markdown": md,
        "images": _extract_image_candidates(url),
        "metadata": {"extracted_by": "browser"},
    }
    if page_text or md:
        return result
    return None


def _try_document_extraction(url: str) -> dict | None:
    path = urlparse(url).path.lower()
    tmp_path = _download_to_temp(url, suffix=Path(path).suffix or ".bin")
    if not tmp_path:
        return None

    try:
        if path.endswith(".pdf"):
            return _extract_pdf(tmp_path)
        if path.endswith(".docx") or path.endswith(".doc"):
            return _extract_docx(tmp_path)
        if path.endswith(".xlsx") or path.endswith(".xls"):
            return _extract_xlsx(tmp_path)
        if path.endswith(".pptx"):
            return _extract_pptx(tmp_path)
    except Exception as exc:
        return {"text": "", "markdown": "", "metadata": {"doc_error": str(exc)}}
    finally:
        try:
            Path(tmp_path).unlink(missing_ok=True)
        except Exception:
            pass
    return None


def _try_metadata_only(url: str) -> dict | None:
    title = ""
    description = ""
    try:
        import httpx
        resp = httpx.get(url, headers={"User-Agent": "VYREN/1.0"}, follow_redirects=True, timeout=8)
        if resp.status_code < 400 and resp.text:
            title = _extract_title(resp.text)
            description = _extract_meta(resp.text, "description")
    except Exception:
        return None
    if title or description:
        return {"title": title, "description": description, "metadata": {"mode": "metadata_only"}}
    return None


def _download_to_temp(url: str, suffix: str = "") -> str | None:
    try:
        import httpx
    except Exception:
        return None
    dest = Path(tempfile.gettempdir()) / f"vyren_dl_{abs(hash(url))}{suffix}"
    try:
        with httpx.stream("GET", url, follow_redirects=True, timeout=20) as resp:
            if resp.status_code >= 400:
                return None
            with open(dest, "wb") as f:
                for chunk in resp.iter_bytes():
                    f.write(chunk)
        return str(dest)
    except Exception:
        return None

# ---------------------------------------------------------------------------
# Document extractors
# ---------------------------------------------------------------------------

def _extract_pdf(path: str) -> dict:
    text = ""
    try:
        import fitz
        doc = fitz.open(path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        text = "\n".join(pages)
    except Exception:
        text = ""
    if not text:
        try:
            from vision.ocr import resolve_backend
            result = resolve_backend("auto").detect_text(path)
            text = result.text or ""
        except Exception:
            text = ""
    return {"text": text, "markdown": text, "metadata": {"format": "pdf"}}


def _extract_docx(path: str) -> dict:
    text = ""
    try:
        from docx import Document
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        text = ""
    return {"text": text, "markdown": text, "metadata": {"format": "docx"}}


def _extract_xlsx(path: str) -> dict:
    text = ""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) for v in row if v is not None]
                if vals:
                    rows.append(" | ".join(vals))
        text = "\n".join(rows)
    except Exception:
        text = ""
    return {"text": text, "markdown": text, "metadata": {"format": "xlsx"}}


def _extract_pptx(path: str) -> dict:
    text = ""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
            if lines:
                slides.append("\n".join(lines))
        text = "\n\n".join(slides)
    except Exception:
        text = ""
    return {"text": text, "markdown": text, "metadata": {"format": "pptx"}}

# ---------------------------------------------------------------------------
# HTML helpers
# ---------------------------------------------------------------------------

def _extract_title(html: str) -> str:
    m = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
    if not m:
        return ""
    title = re.sub(r"<[^>]+>", "", m.group(1))
    return title.strip()


def _extract_meta(html: str, name: str) -> str:
    pattern = re.compile(
        r"<meta[^>]+(?:name|property)=['\"](?:og:)?{}['\"][^>]+content=['\"]([^'\"]+)['\"]".format(re.escape(name)),
        re.IGNORECASE,
    )
    m = pattern.search(html)
    if m:
        return m.group(1).strip()
    pattern2 = re.compile(
        r"<meta[^>]+content=['\"]([^'\"]+)['\"][^>]+(?:name|property)=['\"](?:og:)?{}['\"]".format(re.escape(name)),
        re.IGNORECASE,
    )
    m2 = pattern2.search(html)
    if m2:
        return m2.group(1).strip()
    return ""


def _extract_canonical(html: str, fallback: str) -> str:
    m = re.search(r"<link[^>]+rel=['\"]canonical['\"][^>]+href=['\"]([^'\"]+)['\"]", html, re.IGNORECASE)
    if m:
        return m.group(1).strip()
    return fallback


def _extract_readable_text(html: str) -> str:
    html = re.sub(r"<script[\s\S]*?</script>", " ", html, flags=re.IGNORECASE)
    html = re.sub(r"<style[\s\S]*?</style>", " ", html, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", html)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&lt;", "<", text)
    text = re.sub(r"&gt;", ">", text)
    text = re.sub(r"&amp;", "&", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _html_to_markdown(html: str) -> str:
    text = _extract_readable_text(html)
    if not text:
        return ""
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return "\n\n".join(lines)


def _text_to_markdown(text: str) -> str:
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip() if text else ""


def _extract_image_candidates(base_url: str) -> list[str]:
    return []

# ---------------------------------------------------------------------------
# Resource store
# ---------------------------------------------------------------------------

class UrlMemory:
    """Session-scoped and optionally persistent URL store."""

    def __init__(self, persist: bool = False, storage_path: Path | None = None) -> None:
        self._session: dict[str, UrlResource] = {}
        self._persist = persist
        self._path = storage_path or Path(tempfile.gettempdir()) / "vyren_url_memory.json"
        if persist and self._path.exists():
            try:
                data = json.loads(self._path.read_text(encoding="utf-8"))
                for item in data:
                    r = UrlResource(**item)
                    self._session[r.url] = r
            except Exception:
                pass

    def put(self, resource: UrlResource) -> None:
        self._session[resource.url] = resource
        if self._persist:
            try:
                payload = [r.to_dict() for r in self._session.values()]
                self._path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
            except Exception:
                pass

    def get(self, url: str) -> UrlResource | None:
        return self._session.get(url)

    def recent(self, limit: int = 20) -> list[UrlResource]:
        items = sorted(self._session.values(), key=lambda r: r.retrieved_at, reverse=True)
        return items[: max(1, limit)]

    def all(self) -> list[UrlResource]:
        return list(self._session.values())
